#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""わんさかんさいの週次SEOレポートを作る。

入力は Search Console の日別データを写したCSV（date,clicks,impressions,ctr,position）。
直近の「完全な1週間（月曜〜日曜）」を自動で選び、前週と比べたレポートを
グラフつきのPDFとPPTXで書き出す。

PDFは matplotlib で作る。PowerPointのCOM自動操作は無人実行で次の不具合を連続して
起こしたため使わない（引数省略での型変換エラー／連続実行でのcom_error／
ファイルロック ~$*.pptx の残留／ウィンドウなしゾンビプロセスの残留／
CoUninitializeでのセグメンテーション違反）。
PPTX（python-pptx）は一度も失敗していないので、あとから手で直せるよう併せて出す。

使い方:
    python scripts/seo_report.py --csv daily.csv --out-dir ./out
    python scripts/seo_report.py --csv daily.csv --out-dir ./out --week 2026-08-17
"""
import argparse
import csv
import datetime
import io
import os
import sys

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

# ---- 配色（サイトのアースカラーに合わせる）----
CREAM = "#faf7f2"
INK = "#3f3a33"
MUTE = "#7a7266"
RULE = "#dfd7c9"
GREEN = "#5c8a5c"     # クリック
BROWN = "#b08655"     # 表示回数
GRAY = "#c4bcae"      # 比較用（過去）
ACCENT = "#c2703d"    # 強調

plt.rcParams["font.family"] = "Yu Gothic"
plt.rcParams["axes.unicode_minus"] = False
plt.rcParams["axes.spines.top"] = False
plt.rcParams["axes.spines.right"] = False
plt.rcParams["axes.edgecolor"] = "#b9b0a2"
plt.rcParams["axes.labelcolor"] = "#5b5145"
plt.rcParams["text.color"] = INK
plt.rcParams["xtick.color"] = MUTE
plt.rcParams["ytick.color"] = MUTE
plt.rcParams["grid.color"] = "#e6e0d6"


# =====================================================================
# データ
# =====================================================================
def load(path):
    """日別CSVを読む。日付順に並べ替えて返す。"""
    rows = []
    with io.open(path, encoding="utf-8-sig") as f:
        for r in csv.DictReader(f):
            rows.append({
                "d": datetime.date.fromisoformat(r["date"].strip().replace("/", "-")),
                "c": int(str(r["clicks"]).replace(",", "")),
                "i": int(str(r["impressions"]).replace(",", "")),
                "pos": float(r["position"]),
            })
    if not rows:
        raise SystemExit("CSVが空です")
    rows.sort(key=lambda r: r["d"])
    return rows


def weeks_of(rows):
    """月曜はじまりで週ごとに集計する。7日そろっていない週は捨てる。

    平均掲載順位は表示回数で重み付けする（日ごとの単純平均だと
    表示の少ない日が同じ重みで効いてしまい実態とずれる）。
    """
    acc = {}
    for r in rows:
        wk = r["d"] - datetime.timedelta(days=r["d"].weekday())
        a = acc.setdefault(wk, {"c": 0, "i": 0, "ps": 0.0, "n": 0})
        a["c"] += r["c"]
        a["i"] += r["i"]
        a["ps"] += r["pos"] * r["i"]
        a["n"] += 1
    out = {}
    for k, a in acc.items():
        if a["n"] != 7:
            continue
        out[k] = {
            "c": a["c"], "i": a["i"], "n": a["n"],
            "ctr": (a["c"] / a["i"] * 100) if a["i"] else 0.0,
            "pos": (a["ps"] / a["i"]) if a["i"] else 0.0,
        }
    if not out:
        raise SystemExit("完全な1週間ぶんのデータがありません")
    return out


def moving_avg(vals, win=7):
    return [sum(vals[max(0, i - win + 1):i + 1]) / len(vals[max(0, i - win + 1):i + 1])
            for i in range(len(vals))]


def jp_range(monday):
    """月曜の日付から「8月17日（月）〜8月23日（日）」を作る。"""
    sunday = monday + datetime.timedelta(days=6)
    return f"{monday.month}月{monday.day}日（月）〜{sunday.month}月{sunday.day}日（日）"


# =====================================================================
# グラフ
# =====================================================================
def make_charts(rows, weeks, cur_key, out_dir):
    ch = os.path.join(out_dir, "charts")
    os.makedirs(ch, exist_ok=True)
    paths = {}

    def save(fig, name):
        p = os.path.join(ch, name)
        fig.savefig(p, dpi=170, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        paths[name] = p

    xs = [r["d"] for r in rows]
    span = f"{xs[0].month}月{xs[0].day}日〜{xs[-1].month}月{xs[-1].day}日"

    # クリック数
    cl = [r["c"] for r in rows]
    fig, ax = plt.subplots(figsize=(11, 4.2))
    ax.fill_between(xs, cl, color=GREEN, alpha=0.18)
    ax.plot(xs, cl, color=GREEN, lw=1.0, alpha=0.55, label="1日ごと")
    ax.plot(xs, moving_avg(cl), color=GREEN, lw=2.8, label="7日移動平均")
    ax.set_ylabel("クリック数")
    ax.grid(axis="y", lw=0.8)
    ax.legend(frameon=False, loc="upper left")
    ax.set_title(f"検索からのクリック数の推移（{span}）", loc="left", fontsize=13, pad=12)
    fig.autofmt_xdate()
    save(fig, "clicks.png")

    # 表示回数
    im = [r["i"] for r in rows]
    fig, ax = plt.subplots(figsize=(11, 4.2))
    ax.fill_between(xs, im, color=BROWN, alpha=0.16)
    ax.plot(xs, im, color=BROWN, lw=1.0, alpha=0.55, label="1日ごと")
    ax.plot(xs, moving_avg(im), color=BROWN, lw=2.8, label="7日移動平均")
    ax.set_ylabel("表示回数")
    ax.yaxis.set_major_formatter(FuncFormatter(lambda v, p: f"{int(v):,}"))
    ax.grid(axis="y", lw=0.8)
    ax.legend(frameon=False, loc="upper left")
    ax.set_title("検索結果に表示された回数の推移", loc="left", fontsize=13, pad=12)
    fig.autofmt_xdate()
    save(fig, "impressions.png")

    # 週ごと（直近12週）
    keys = sorted(k for k in weeks if k <= cur_key)[-12:]
    labels = [f"{k.month}/{k.day}" for k in keys]
    vals = [weeks[k]["c"] for k in keys]
    colors = [ACCENT if k == cur_key else GRAY for k in keys]
    fig, ax = plt.subplots(figsize=(11, 4.2))
    b = ax.bar(labels, vals, color=colors, width=0.68)
    ax.bar_label(b, labels=[f"{v:,}" for v in vals], padding=3, fontsize=9)
    ax.set_ylabel("クリック数")
    ax.set_xlabel("週のはじまり（月曜）")
    ax.set_ylim(0, max(vals) * 1.18)
    ax.grid(axis="y", lw=0.8)
    ax.set_title("週ごとのクリック数（月曜〜日曜。オレンジが対象週）",
                 loc="left", fontsize=13, pad=12)
    save(fig, "weekly.png")

    # 掲載順位
    ps = [r["pos"] for r in rows]
    fig, ax = plt.subplots(figsize=(11, 3.9))
    ax.plot(xs, ps, color="#a9a093", lw=0.9, alpha=0.6, label="1日ごと")
    ax.plot(xs, moving_avg(ps), color=ACCENT, lw=2.8, label="7日移動平均")
    ax.invert_yaxis()
    ax.set_ylabel("平均掲載順位")
    ax.grid(axis="y", lw=0.8)
    ax.legend(frameon=False, loc="lower left")
    ax.set_title("平均掲載順位の推移（上にいくほど順位が良い）",
                 loc="left", fontsize=13, pad=12)
    fig.autofmt_xdate()
    save(fig, "position.png")

    return paths


# =====================================================================
# 文章
# =====================================================================
def build_meta(rows, weeks, cur_key, prev_key):
    cur, prev = weeks[cur_key], weeks[prev_key]
    dc = (cur["c"] - prev["c"]) / prev["c"] * 100 if prev["c"] else 0.0
    di = (cur["i"] - prev["i"]) / prev["i"] * 100 if prev["i"] else 0.0
    dp = prev["pos"] - cur["pos"]            # プラスなら順位が上がった

    # 週単位は上下しやすいので、直近12週の中での位置づけを必ず添える
    keys12 = sorted(k for k in weeks if k <= cur_key)[-12:]
    vals12 = [weeks[k]["c"] for k in keys12]
    n12 = len(vals12)
    rank = sorted(vals12, reverse=True).index(cur["c"]) + 1
    if rank <= max(3, n12 // 4):
        context = f"ただし直近{n12}週では{rank}番目に多い週で、水準としては高い方です。"
    elif rank >= n12 - 2:
        context = f"直近{n12}週では下から{n12 - rank + 1}番目で、水準としても低めです。要注意。"
    else:
        context = f"直近{n12}週では{rank}番目で、水準としては真ん中あたりです。"

    if dc >= 10:
        headline = f"クリックが前週より{dc:.0f}%　伸びました"
    elif dc <= -10:
        headline = f"クリックが前週より{abs(dc):.0f}%　減りました"
    else:
        headline = "クリックはほぼ横ばいです"

    avg12 = sum(vals12) / n12
    lead = (f"対象は{jp_range(cur_key)}の1週間です。"
            f"クリック{cur['c']:,}回、表示{cur['i']:,}回。"
            f"直近{n12}週の平均は{avg12:,.0f}回です。")

    reading = (
        f"クリックは{prev['c']:,} → {cur['c']:,}（{dc:+.0f}%）、"
        f"表示は{prev['i']:,} → {cur['i']:,}（{di:+.0f}%）。\n"
        f"{context}\n"
        f"平均掲載順位は{prev['pos']:.1f}位 → {cur['pos']:.1f}位"
        f"（{'改善' if dp > 0.05 else '悪化' if dp < -0.05 else 'ほぼ変化なし'}）。\n"
        "1週だけの増減はニュースや天気でも動きます。次ページの12週の並びで見てください。"
    )

    weekly_note = (f"対象週は{cur['c']:,}回。直近{n12}週の平均{avg12:,.0f}回に対して"
                   f"{cur['c'] / avg12 * 100 - 100:+.0f}%、{rank}番目の水準です。")

    tot_c = sum(r["c"] for r in rows)
    tot_i = sum(r["i"] for r in rows)
    first, last = rows[0]["d"], rows[-1]["d"]
    clicks_note = (f"データのある{len(rows)}日間で合計{tot_c:,}クリック。"
                   f"1日あたり平均{tot_c / len(rows):.0f}回です。")
    imp_note = (f"同じ期間で合計{tot_i:,}回表示され、全体のクリック率は{tot_c / tot_i * 100:.1f}%。"
                "表示が増えていればクリックは後からついてきます。")
    pos_note = (f"対象週の平均掲載順位は{cur['pos']:.1f}位。"
                "10位以内なら検索結果の1ページ目に入っている状態です。")

    summary = (
        f"■ 対象週：{jp_range(cur_key)}\n\n"
        f"・クリック {cur['c']:,}回（前週比 {dc:+.0f}%）\n"
        f"・表示 {cur['i']:,}回（前週比 {di:+.0f}%）\n"
        f"・クリック率 {cur['ctr']:.1f}%　平均掲載順位 {cur['pos']:.1f}位\n\n"
        f"{context}\n\n"
        f"集計期間全体（{first.month}/{first.day}〜{last.month}/{last.day}・{len(rows)}日）では"
        f"合計{tot_c:,}クリック、{tot_i:,}回表示されています。"
    )

    return {
        "cur": cur, "prev": prev,
        "cur_key": cur_key, "prev_key": prev_key,
        "subtitle": f"対象週　{jp_range(cur_key)}（前週と比較）",
        "headline": headline, "lead": lead, "reading": reading,
        "weekly_note": weekly_note, "clicks_note": clicks_note,
        "imp_note": imp_note, "pos_note": pos_note, "summary": summary,
        "dc": dc, "di": di, "rank": rank, "n12": n12,
    }


# =====================================================================
# PDF（matplotlib・Office非依存）
# =====================================================================
def build_pdf(charts, meta, path):
    """1ページ = 1つのFigure。16:9（13.333x7.5インチ）で作る。"""
    from matplotlib.backends.backend_pdf import PdfPages
    import matplotlib.image as mpimg

    W, H = 13.333, 7.5

    def page(bg=CREAM):
        fig = plt.figure(figsize=(W, H), dpi=150)
        fig.patch.set_facecolor(bg)
        return fig

    def T(fig, s, x, y, size=14, color=INK, weight="normal"):
        fig.text(x / W, 1 - y / H, s, fontsize=size, color=color,
                 fontweight=weight, va="top", ha="left", linespacing=1.75)

    def HR(fig, x, y, w):
        fig.add_artist(plt.Line2D([x / W, (x + w) / W], [1 - y / H, 1 - y / H],
                                  color=RULE, lw=1.2, transform=fig.transFigure))

    def head(fig, title, sub=None):
        T(fig, title, 0.7, 0.5, size=24, weight="bold")
        if sub:
            T(fig, sub, 0.7, 1.15, size=11.5, color=MUTE)
        HR(fig, 0.7, 1.6, 11.9)

    def img(fig, key, x, y, w, max_h=None):
        """画像を置いて下端（インチ）を返す。実寸から高さを出して文字と重ならないようにする。"""
        p = charts.get(key)
        if not p or not os.path.exists(p):
            return y
        im = mpimg.imread(p)
        ih, iw = im.shape[0], im.shape[1]
        h = w * ih / iw
        if max_h and h > max_h:
            h = max_h
            w = h * iw / ih
            x = 0.75 + (11.5 - w) / 2
        ax = fig.add_axes([x / W, 1 - (y + h) / H, w / W, h / H])
        ax.imshow(im)
        ax.axis("off")
        return y + h

    def kpi(fig, x, y, w, label, value, sub, vc=INK):
        ax = fig.add_axes([x / W, 1 - (y + 1.7) / H, w / W, 1.7 / H])
        ax.set_xticks([])
        ax.set_yticks([])
        ax.set_facecolor("white")
        for sp in ax.spines.values():
            sp.set_color(RULE)
        ax.text(0.07, 0.84, label, fontsize=10.5, color=MUTE, va="top",
                transform=ax.transAxes)
        ax.text(0.07, 0.55, value, fontsize=25, color=vc, fontweight="bold",
                va="center", transform=ax.transAxes)
        ax.text(0.07, 0.13, sub, fontsize=9, color=MUTE, va="center",
                transform=ax.transAxes)

    cur, prev = meta["cur"], meta["prev"]
    dc, di = meta["dc"], meta["di"]
    cc = ACCENT if dc < 0 else GREEN
    ctr_c = ACCENT if cur["ctr"] < prev["ctr"] - 0.05 else GREEN

    with PdfPages(path) as pdf:
        # 表紙
        fig = page("white")
        fig.add_artist(plt.Rectangle((0, 1 - 2.7 / H), 1, 2.7 / H,
                                     transform=fig.transFigure, color=GREEN, zorder=0))
        T(fig, "わんさかんさい　週次レポート", 0.9, 0.95, size=30, color="white", weight="bold")
        T(fig, meta["subtitle"], 0.9, 1.9, size=12.5, color="#e6efe6")
        T(fig, meta["headline"], 0.9, 3.4, size=19, color=cc, weight="bold")
        T(fig, meta["lead"], 0.9, 4.2, size=12, color=MUTE)
        pdf.savefig(fig, facecolor=fig.get_facecolor())
        plt.close(fig)

        # サマリー
        fig = page()
        head(fig, "対象週のサマリー", meta["subtitle"])
        kpi(fig, 0.70, 2.05, 2.85, "クリック数", f"{cur['c']:,}",
            f"前週 {prev['c']:,} → {dc:+.0f}%", cc)
        kpi(fig, 3.75, 2.05, 2.85, "表示回数", f"{cur['i']:,}",
            f"前週 {prev['i']:,} → {di:+.0f}%", cc)
        kpi(fig, 6.80, 2.05, 2.85, "クリック率", f"{cur['ctr']:.1f}%",
            f"前週 {prev['ctr']:.1f}%", ctr_c)
        kpi(fig, 9.85, 2.05, 2.78, "平均掲載順位", f"{cur['pos']:.1f}",
            f"前週 {prev['pos']:.1f}", INK)
        T(fig, "読み解き", 0.7, 4.3, size=14.5, color=GREEN, weight="bold")
        T(fig, meta["reading"], 0.7, 4.8, size=12)
        pdf.savefig(fig, facecolor=fig.get_facecolor())
        plt.close(fig)

        # グラフ各ページ
        for key, title, sub, note in [
            ("weekly.png", "直近12週の推移",
             "1週だけを見ると季節や天気のブレに振り回されるため、必ず並べて確認します",
             meta["weekly_note"]),
            ("clicks.png", "クリック数の推移",
             "細い線が1日ごと、太い線が7日移動平均", meta["clicks_note"]),
            ("impressions.png", "検索結果に表示された回数",
             "クリックの手前にある指標", meta["imp_note"]),
            ("position.png", "平均掲載順位の推移",
             "グラフは上にいくほど順位が良い状態です", meta["pos_note"]),
        ]:
            fig = page()
            head(fig, title, sub)
            b = img(fig, key, 0.75, 1.95, 11.5, max_h=4.3)
            T(fig, note, 0.75, b + 0.32, size=12.5, color=GREEN, weight="bold")
            pdf.savefig(fig, facecolor=fig.get_facecolor())
            plt.close(fig)

        # まとめ
        fig = page()
        head(fig, "まとめ")
        T(fig, meta["summary"], 0.9, 2.15, size=13)
        HR(fig, 0.9, 6.05, 11.5)
        T(fig, "数字はGoogle Search Consoleの実測値です。データは2〜3日遅れて確定します。",
          0.9, 6.35, size=11, color=MUTE)
        pdf.savefig(fig, facecolor=fig.get_facecolor())
        plt.close(fig)

    return path


# =====================================================================
# PPTX（あとから手で直す用）
# =====================================================================
def build_pptx(charts, meta, path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[pptx] python-pptx が無いのでPPTXは省略します")
        return None
    from PIL import Image

    def rgb(h):
        return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide(bg=CREAM):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = rgb(bg)
        return s

    def text(s, t, x, y, w, h, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(str(t).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = "Yu Gothic"
            r.font.color.rgb = rgb(color)
        return tb

    def pic(s, key, x, y, w, max_h=None):
        p = charts.get(key)
        if not p or not os.path.exists(p):
            return y
        iw, ih = Image.open(p).size
        h = w * ih / iw
        if max_h and h > max_h:
            h = max_h
            w2 = h * iw / ih
            x = x + (w - w2) / 2
            w = w2
        s.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(w))
        return y + h

    def head(s, title, sub=None):
        text(s, title, 0.7, 0.35, 11.9, 0.8, size=26, bold=True)
        if sub:
            text(s, sub, 0.7, 1.05, 11.9, 0.5, size=12, color=MUTE)

    cur, prev = meta["cur"], meta["prev"]
    dc, di = meta["dc"], meta["di"]
    cc = ACCENT if dc < 0 else GREEN
    ctr_c = ACCENT if cur["ctr"] < prev["ctr"] - 0.05 else GREEN

    # 表紙
    s = slide("#ffffff")
    band = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(2.7))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(GREEN)
    band.line.fill.background()
    band.shadow.inherit = False
    text(s, "わんさかんさい　週次レポート", 0.85, 0.75, 11.6, 1.0, size=32, bold=True, color="#ffffff")
    text(s, meta["subtitle"], 0.85, 1.75, 11.6, 0.5, size=13, color="#e6efe6")
    text(s, meta["headline"], 0.85, 3.2, 11.6, 0.8, size=20, bold=True, color=cc)
    text(s, meta["lead"], 0.85, 4.05, 11.6, 1.2, size=12.5, color=MUTE)

    # サマリー
    s = slide()
    head(s, "対象週のサマリー", meta["subtitle"])
    for i, (label, value, sub, col) in enumerate([
        ("クリック数", f"{cur['c']:,}", f"前週 {prev['c']:,} → {dc:+.0f}%", cc),
        ("表示回数", f"{cur['i']:,}", f"前週 {prev['i']:,} → {di:+.0f}%", cc),
        ("クリック率", f"{cur['ctr']:.1f}%", f"前週 {prev['ctr']:.1f}%", ctr_c),
        ("平均掲載順位", f"{cur['pos']:.1f}", f"前週 {prev['pos']:.1f}", INK),
    ]):
        x = 0.7 + i * 3.05
        card = s.shapes.add_shape(1, Inches(x), Inches(1.95), Inches(2.85), Inches(1.7))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("#ffffff")
        card.line.color.rgb = rgb(RULE)
        card.shadow.inherit = False
        text(s, label, x + 0.2, 2.05, 2.5, 0.4, size=11, color=MUTE)
        text(s, value, x + 0.2, 2.4, 2.5, 0.7, size=26, bold=True, color=col)
        text(s, sub, x + 0.2, 3.15, 2.5, 0.4, size=9.5, color=MUTE)
    text(s, "読み解き", 0.7, 4.2, 11.9, 0.4, size=15, bold=True, color=GREEN)
    text(s, meta["reading"], 0.7, 4.7, 11.9, 2.3, size=12.5)

    # グラフ各ページ
    for key, title, sub, note in [
        ("weekly.png", "直近12週の推移",
         "1週だけを見ると季節や天気のブレに振り回されるため、必ず並べて確認します",
         meta["weekly_note"]),
        ("clicks.png", "クリック数の推移",
         "細い線が1日ごと、太い線が7日移動平均", meta["clicks_note"]),
        ("impressions.png", "検索結果に表示された回数",
         "クリックの手前にある指標", meta["imp_note"]),
        ("position.png", "平均掲載順位の推移",
         "グラフは上にいくほど順位が良い状態です", meta["pos_note"]),
    ]:
        s = slide()
        head(s, title, sub)
        b = pic(s, key, 0.75, 1.9, 11.5, max_h=4.3)
        text(s, note, 0.75, b + 0.25, 11.8, 0.9, size=14, bold=True, color=GREEN)

    # まとめ
    s = slide()
    head(s, "まとめ")
    text(s, meta["summary"], 0.9, 1.9, 11.5, 4.2, size=13.5)
    text(s, "数字はGoogle Search Consoleの実測値です。データは2〜3日遅れて確定します。",
         0.9, 6.4, 11.5, 0.5, size=11.5, color=MUTE)

    prs.save(path)
    return path


# =====================================================================
def main():
    # 実行元のシェルに関係なく出力をUTF-8にそろえる（自動タスクがログに取り込むため）
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    ap = argparse.ArgumentParser()
    ap.add_argument("--csv", required=True, help="日別データのCSV")
    ap.add_argument("--out-dir", required=True, help="出力先フォルダ")
    ap.add_argument("--week", default=None,
                    help="対象週の月曜（YYYY-MM-DD）。省略時は直近の完全な週")
    ap.add_argument("--no-pptx", action="store_true")
    ap.add_argument("--no-pdf", action="store_true")
    a = ap.parse_args()

    rows = load(a.csv)
    weeks = weeks_of(rows)
    keys = sorted(weeks)

    if a.week:
        cur_key = datetime.date.fromisoformat(a.week)
        if cur_key not in weeks:
            raise SystemExit(f"{a.week} の週は完全な7日ぶんがありません")
    else:
        cur_key = keys[-1]
    prev_key = cur_key - datetime.timedelta(days=7)
    if prev_key not in weeks:
        raise SystemExit("前週のデータがそろっていないので比較できません")

    os.makedirs(a.out_dir, exist_ok=True)
    charts = make_charts(rows, weeks, cur_key, a.out_dir)
    meta = build_meta(rows, weeks, cur_key, prev_key)

    stamp = cur_key.isoformat()
    pptx = pdf = None
    if not a.no_pptx:
        pptx = build_pptx(charts, meta, os.path.join(a.out_dir, f"wansakansai_weekly_{stamp}.pptx"))
        if pptx:
            print(f"PPTX: {pptx}")
    if not a.no_pdf:
        pdf = build_pdf(charts, meta, os.path.join(a.out_dir, f"wansakansai_weekly_{stamp}.pdf"))
        print(f"PDF : {pdf}")

    cur = meta["cur"]
    # 自動タスクがメール本文に使うための1行
    print(f"[key] 期間={cur_key.month:02d}/{cur_key.day:02d}〜"
          f"{(cur_key + datetime.timedelta(days=6)).month:02d}/"
          f"{(cur_key + datetime.timedelta(days=6)).day:02d} "
          f"クリック={cur['c']} 前週={meta['prev']['c']} 増減={meta['dc']:+.0f}% "
          f"表示={cur['i']} CTR={cur['ctr']:.1f}% 順位={cur['pos']:.1f} "
          f"順位付け={meta['rank']}/{meta['n12']}週")
    return 0


if __name__ == "__main__":
    sys.exit(main())
